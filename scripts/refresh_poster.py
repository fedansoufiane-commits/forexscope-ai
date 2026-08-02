"""Refresh the A3 poster from the versioned diagnostics artifacts.

The poster is the one deliverable that carried hand-typed metrics and exported
chart images, so it silently kept the pre-1.0 numbers (Accuracy 55,3 %, ROC-AUC
0,588, 5-Fold Stratified CV) long after the purged out-of-time rebuild replaced
them. This script rewrites every number and both result figures straight from
`models/diagnostics.json` and `models/learning_curve.json`, so the poster can
only ever show what the trainer actually produced.

Run after `scripts/train_and_diagnose.py`:

    python3 scripts/refresh_poster.py

The correlation matrix (Abb. 1) is model-independent and stays as it is.
"""
from __future__ import annotations

import copy
import json
import sys
from pathlib import Path

ROOT = Path(__file__).resolve().parent.parent
sys.path.insert(0, str(ROOT))

from pptx import Presentation  # noqa: E402
from pptx.util import Emu  # noqa: E402

POSTER = ROOT / "Präsentation" / "WealthScope_Poster_A3.pptx"
DIAGNOSTICS = ROOT / "models" / "diagnostics.json"
LEARNING_CURVE = ROOT / "models" / "learning_curve.json"

# Shape indices on the single poster slide, resolved once by inspection and
# re-verified below against the text they are expected to contain.
IDX_ACCURACY = 37
IDX_AUC = 40
IDX_AUC_LABEL = 41
IDX_CONFUSION_PIC = 47
IDX_LEARNING_PIC = 50
IDX_LEARNING_CAPTION = 51
IDX_VERDICT = 53
IDX_METHOD = 56

RENDER_SCALE = 2  # plotly px -> ~2048 px wide, matching the existing figures


def de(value: str) -> str:
    return value.replace(".", ",")


def render_figures(out_dir: Path) -> tuple[Path, Path]:
    """Render the two result charts with the app's own chart code."""
    import plotly.io as pio

    from src.diagnostics import chart_confusion_matrix, chart_learning_curve

    diag = json.loads(DIAGNOSTICS.read_text(encoding="utf-8"))
    lc = json.loads(LEARNING_CURVE.read_text(encoding="utf-8"))

    out_dir.mkdir(parents=True, exist_ok=True)
    cm_path = out_dir / "poster_confusion_matrix.png"
    lc_path = out_dir / "poster_learning_curve.png"

    cm_fig = chart_confusion_matrix(diag, "Hell")
    cm_fig.update_layout(width=1024, height=402)
    pio.write_image(cm_fig, cm_path, scale=RENDER_SCALE)

    lc_fig = chart_learning_curve(lc, "Hell")
    lc_fig.update_layout(width=1024, height=442)
    pio.write_image(lc_fig, lc_path, scale=RENDER_SCALE)

    return cm_path, lc_path


def set_text(shape, text: str) -> None:
    """Replace the text while keeping the first run's formatting."""
    para = shape.text_frame.paragraphs[0]
    for run in list(para.runs)[1:]:
        run._r.getparent().remove(run._r)
    para.runs[0].text = text


def set_lead_and_body(shape, lead: str, body: str) -> None:
    """Two-run paragraph: bold lead sentence followed by the body text."""
    para = shape.text_frame.paragraphs[0]
    runs = list(para.runs)
    for run in runs[2:]:
        run._r.getparent().remove(run._r)
    runs[0].text = lead
    runs[1].text = body


def replace_picture(slide, index: int, image_path: Path):
    """Swap a picture in place, preserving position, size and z-order."""
    old = slide.shapes[index]
    left, top, width, height = old.left, old.top, old.width, old.height
    new = slide.shapes.add_picture(str(image_path), left, top, width, height)
    old_el, new_el = old._element, new._element
    old_el.addprevious(new_el)
    old_el.getparent().remove(old_el)
    return new


def main() -> int:
    diag = json.loads(DIAGNOSTICS.read_text(encoding="utf-8"))
    lc = json.loads(LEARNING_CURVE.read_text(encoding="utf-8"))
    m = diag["test_metrics"]
    val = diag["validation"]

    accuracy = de(f"{m['accuracy'] * 100:.1f} %")
    auc = de(f"{m['roc_auc']:.3f}")
    baseline = de(f"{m['majority_baseline'] * 100:.1f} %")
    folds = lc["cv_n_splits"]
    test_from = val["test_start"][:4]
    test_to = val["test_end"][:4]
    train_to = val["train_end"][:4]
    purge = val["purge_trading_days"]

    experiments = json.loads(
        (ROOT / "models" / "validation_experiments.json").read_text(encoding="utf-8")
    )
    summary = experiments["split_comparison_summary"]
    naive_auc = de(f"{summary['leaky_roc_auc']:.3f}")
    factor = de(f"{summary['apparent_signal_factor']:.1f}")

    prs = Presentation(POSTER)
    slide = prs.slides[0]
    shapes = slide.shapes

    # Fail loudly if the deck was re-laid-out and the indices moved.
    expected = {
        IDX_ACCURACY: "%",
        IDX_AUC: "0,",
        IDX_AUC_LABEL: "ROC-AUC",
        IDX_LEARNING_CAPTION: "Abb. 3",
        IDX_VERDICT: "Warum",
        IDX_METHOD: "RandomForest",
    }
    for idx, needle in expected.items():
        actual = shapes[idx].text_frame.text
        if needle not in actual:
            raise SystemExit(
                f"Shape {idx} contains {actual[:60]!r}, expected to find {needle!r}. "
                "The poster layout changed - re-resolve the shape indices."
            )
    for idx in (IDX_CONFUSION_PIC, IDX_LEARNING_PIC):
        if shapes[idx].shape_type != 13:
            raise SystemExit(f"Shape {idx} is not a picture; re-resolve the indices.")

    scratch = ROOT / "models" / ".poster_figures"
    cm_path, lc_path = render_figures(scratch)

    set_text(shapes[IDX_ACCURACY], accuracy)
    set_text(shapes[IDX_AUC], auc)
    set_text(shapes[IDX_AUC_LABEL], "ROC-AUC (Out-of-Time)")
    set_text(
        shapes[IDX_LEARNING_CAPTION],
        f"Abb. 3: Lernkurve (ROC-AUC vs. Trainingsgröße, {folds} expandierende "
        "Walk-forward-Folds) — Quelle: eigene Berechnung",
    )
    set_lead_and_body(
        shapes[IDX_VERDICT],
        f"Warum {accuracy} das ehrliche Ergebnis ist.  ",
        f"Der Test ist zeitlich unangetastet: Training bis {train_to}, Prüfung auf "
        f"{test_from}–{test_to}, dazwischen {purge} Handelstage Sperrzone. So bleibt die "
        f"Accuracy unter der Mehrheitsbaseline von {baseline} und der ROC-AUC bei {auc} — "
        "H1 ist damit falsifiziert, nicht ungeprüft. Genau das erwartet die "
        f"Effizienzmarkthypothese (Fama 1970). Ein naiver Zufalls-Split bewertet dieselben "
        f"Daten mit {naive_auc}: ein {factor}-mal größeres scheinbares Signal.",
    )
    set_text(
        shapes[IDX_METHOD],
        "RandomForest (Median-Imputer → RF; StandardScaler nur bei LogReg und Linear "
        f"SVM, Bäume sind skaleninvariant). Purged Out-of-Time-Test plus {folds} "
        "expandierende Walk-forward-Folds, alles nur im Trainingsfenster gefittet. "
        "8 Features aus OHLCV, Ziel target_20d.",
    )

    replace_picture(slide, IDX_LEARNING_PIC, lc_path)
    replace_picture(slide, IDX_CONFUSION_PIC, cm_path)

    prs.save(POSTER)
    print(f"updated: {POSTER}")
    print(f"  Accuracy       {accuracy}  (Baseline {baseline})")
    print(f"  ROC-AUC        {auc}")
    print(f"  Lernkurve      {folds} Walk-forward-Folds")
    print(f"  naiver Split   {naive_auc}  -> Faktor {factor}")
    print("Convert to PDF afterwards, e.g. with:")
    print("  soffice --headless --convert-to pdf --outdir 'Präsentation' "
          "'Präsentation/WealthScope_Poster_A3.pptx'")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
